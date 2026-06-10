import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize presentation
prs = Presentation()
# Set widescreen 16:9 dimensions
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette (Premium ESG theme - harmonized dark/light slate and emerald green)
BG_DARK = RGBColor(18, 30, 49)       # Very dark navy/slate
BG_LIGHT = RGBColor(245, 247, 250)   # Clean light slate gray
PRIMARY = RGBColor(24, 43, 73)       # Deep slate blue
SECONDARY = RGBColor(34, 139, 115)   # Emerald green (ESG theme)
ACCENT = RGBColor(224, 153, 36)      # Amber gold (highlight)
WHITE = RGBColor(255, 255, 255)
CARD_BG = RGBColor(255, 255, 255)
TEXT_DARK = RGBColor(33, 37, 41)
TEXT_LIGHT = RGBColor(240, 240, 240)
TEXT_MUTED = RGBColor(108, 117, 125)

# Fonts
FONT_NAME = "微軟正黑體"
FONT_ENG = "Arial"

# Helper function to remove shape borders safely
def remove_border(shape):
    shape.line.color.rgb = shape.fill.fore_color.rgb

# Helper function to set solid background color
def set_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

# Helper function to create content slide header
def add_slide_header(slide, title_text, dark_mode=False):
    # Add title text box
    tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = FONT_NAME
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT if dark_mode else PRIMARY
    
    # Add accent line under header
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = SECONDARY
    remove_border(line)

# Helper function to create styled cards
def create_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=None, border_width=Pt(1)):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = border_width
    else:
        remove_border(card)
    return card

# Helper function to create card with left accent stripe
def create_elegant_card(slide, left, top, width, height, stripe_color=SECONDARY, bg_color=CARD_BG):
    create_card(slide, left, top, width, height, bg_color)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.12), height)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = stripe_color
    remove_border(stripe)

# Helper to add single bullet/paragraph with styling
def add_bullet_point(tf, bold_prefix, text_content, font_size=15, is_bullet=True, text_color=TEXT_DARK, bold_color=PRIMARY):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
        
    if is_bullet:
        p.level = 0
        p.space_after = Pt(8)
    else:
        p.level = 0
        p.space_after = Pt(6)
        
    p.font.name = FONT_NAME
    p.font.size = Pt(font_size)
    
    # Bold prefix
    if bold_prefix:
        run_bold = p.add_run()
        run_bold.text = bold_prefix
        run_bold.font.bold = True
        run_bold.font.color.rgb = bold_color
    
    # Regular text
    if text_content:
        run_text = p.add_run()
        run_text.text = text_content
        run_text.font.bold = False
        run_text.font.color.rgb = text_color

# Helper to add a standardized card with bullets
def add_card_text_content(slide, left, top, width, height, title, bullets, stripe_color=SECONDARY, font_size=15):
    create_elegant_card(slide, left, top, width, height, stripe_color)
    tb_left = left + Inches(0.4)
    tb_top = top + Inches(0.3)
    tb_width = width - Inches(0.6)
    tb_height = height - Inches(0.6)
    
    tx_box = slide.shapes.add_textbox(tb_left, tb_top, tb_width, tb_height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_NAME
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = PRIMARY
    p_title.space_after = Pt(14)
    
    for bold_prefix, regular_text in bullets:
        add_bullet_point(tf, bold_prefix, regular_text, font_size=font_size, is_bullet=True)

# Helper to add picture or fallback box safely
def add_slide_picture(slide, img_name, left, top, width, height, title_fallback=""):
    script_dir = os.path.dirname(os.path.abspath(__file__))
    img_path = os.path.join(script_dir, img_name)
    if os.path.exists(img_path):
        try:
            slide.shapes.add_picture(img_path, left, top, width, height)
        except Exception as e:
            print(f"Error adding picture {img_name}: {e}")
            create_fallback_box(slide, left, top, width, height, f"[圖片載入失敗: {img_name}]")
    else:
        if os.path.exists(img_name):
            try:
                slide.shapes.add_picture(img_name, left, top, width, height)
                return
            except Exception as e:
                pass
        create_fallback_box(slide, left, top, width, height, f"[圖片未尋獲]\n{title_fallback}")

def create_fallback_box(slide, left, top, width, height, text):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(230, 235, 245)
    box.line.color.rgb = SECONDARY
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_NAME
    p.font.size = Pt(14)
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER

# ------------------------------------------------------------
# Slide 1: Cover Slide (Dark Theme)
# ------------------------------------------------------------
slide_layout = prs.slide_layouts[6] # Blank layout
slide1 = prs.slides.add_slide(slide_layout)
set_background(slide1, BG_DARK)

# Background Accent Shapes
acc1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5), Inches(0), Inches(3.833), Inches(7.5))
acc1.fill.solid()
acc1.fill.fore_color.rgb = RGBColor(24, 38, 60)
remove_border(acc1)

acc2 = slide1.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(7.5), Inches(0), Inches(2.0), Inches(7.5))
acc2.fill.solid()
acc2.fill.fore_color.rgb = RGBColor(24, 38, 60)
acc2.rotation = 180
remove_border(acc2)

acc3 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(13.133), Inches(0), Inches(0.2), Inches(7.5))
acc3.fill.solid()
acc3.fill.fore_color.rgb = SECONDARY
remove_border(acc3)

# Title Box
title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(8.0), Inches(2.5))
tf = title_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

p1 = tf.paragraphs[0]
p1.text = "上市公司 ESG 承諾與證據驗證系統"
p1.font.name = FONT_NAME
p1.font.size = Pt(36)
p1.font.bold = True
p1.font.color.rgb = WHITE

p2 = tf.add_paragraph()
p2.text = "讓 AI 幫我們查核「減碳與永續承諾」是否漂綠"
p2.font.name = FONT_NAME
p2.font.size = Pt(22)
p2.font.bold = True
p2.font.color.rgb = ACCENT
p2.space_before = Pt(12)

p3 = tf.add_paragraph()
p3.text = "Multi-Task ESG Promise and Evidence Verification System"
p3.font.name = FONT_ENG
p3.font.size = Pt(14)
p3.font.color.rgb = SECONDARY
p3.space_before = Pt(12)

# Presenter info
info_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(8.0), Inches(1.5))
tf_info = info_box.text_frame
tf_info.word_wrap = True
tf_info.margin_left = tf_info.margin_right = tf_info.margin_top = tf_info.margin_bottom = 0
pinfo = tf_info.paragraphs[0]
pinfo.text = "2026 AICUP 期末專題報告 (簡明版)"
pinfo.font.name = FONT_NAME
pinfo.font.size = Pt(16)
pinfo.font.color.rgb = TEXT_LIGHT

pinfo2 = tf_info.add_paragraph()
pinfo2.text = "報告人：[您的姓名]      |      報告日期：2026 年 6 月 11 日"
pinfo2.font.name = FONT_NAME
pinfo2.font.size = Pt(13)
pinfo2.font.color.rgb = TEXT_MUTED
pinfo2.space_before = Pt(8)

# ------------------------------------------------------------
# Slide 2: 簡報大綱 (Outline)
# ------------------------------------------------------------
slide2 = prs.slides.add_slide(slide_layout)
set_background(slide2, BG_LIGHT)
add_slide_header(slide2, "簡報大綱 (Outline)")

# Create 4 Cards for Outline Items
card_w = Inches(2.7)
card_h = Inches(4.5)
card_top = Inches(1.8)
card_gap = Inches(0.3)
left_start = Inches(0.8)

outline_items = [
    ("01 動機與目的", [
        ("漂綠質疑：", "很多減碳承諾缺乏實質證據。"),
        ("AI 查核點：", "承諾、時程、證據、品質4維度。")
    ], SECONDARY),
    ("02 學習原理", [
        ("BERT 讀書：", "預訓練打底，再進行微調考卷。"),
        ("多任務大腦：", "一個 AI 大腦回答四個問題。")
    ], PRIMARY),
    ("03 研究與優化", [
        ("左截斷：", "保留文章後半部關鍵證據。"),
        ("特徵二合一：", "融合 CLS 與 Mean Pooling。"),
        ("投票與減重：", "三個臭皮匠與 4GB 顯卡跑訓練。")
    ], SECONDARY),
    ("04 實驗與結論", [
        ("實戰大車拼：", "分數相比 Baseline 大幅提升 37%。"),
        ("訓練與未來：", "收斂穩定，下一步朝本土與 LLM 邁進。")
    ], PRIMARY)
]

for i, (title, bullets, color) in enumerate(outline_items):
    x_pos = left_start + i * (card_w + card_gap)
    create_elegant_card(slide2, x_pos, card_top, card_w, card_h, color)
    
    tb = slide2.shapes.add_textbox(x_pos + Inches(0.2), card_top + Inches(0.3), card_w - Inches(0.3), card_h - Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_NAME
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = PRIMARY
    p_title.space_after = Pt(14)
    
    for prefix, desc in bullets:
        add_bullet_point(tf, prefix, desc, font_size=14, is_bullet=True)

# ------------------------------------------------------------
# Slide 3: 動機與目的 (Motivation & Purpose)
# ------------------------------------------------------------
slide3 = prs.slides.add_slide(slide_layout)
set_background(slide3, BG_LIGHT)
add_slide_header(slide3, "01 為什麼要用 AI 查核 ESG 報告？")

bullets3_left = [
    ("• 永續承諾氾濫：", "許多企業在永續報告書中寫出宏偉的環保與減碳大夢，但往往缺乏實際數字佐證，引起大眾對「漂綠 (Greenwashing)」的質疑。"),
    ("• 人工查核太累：", "報告書動輒數百頁，字多又複雜。仰賴人工一頁頁核對，不僅極度耗時費力，更難以建立客觀、標準化的審查。")
]
bullets3_right = [
    ("• 自動驗證目的：", "開發自動化中文 AI 審查系統，能從永續報告中自動辨識並快速查驗「減碳承諾」與其對應的「佐證數據品質」。"),
    ("• 提升查核效率：", "藉由自動化分類，幫助分析師與稽核員快速篩選出高漂綠風險或低證據品質的企業報告。")
]
add_card_text_content(slide3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "傳統稽核與漂綠挑戰", bullets3_left, SECONDARY)
add_card_text_content(slide3, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), "自動查核目的", bullets3_right, PRIMARY)

# ------------------------------------------------------------
# Slide 4: AI 幫我們看的 4 個 ESG 檢查點 (Tasks)
# ------------------------------------------------------------
slide4 = prs.slides.add_slide(slide_layout)
set_background(slide4, BG_LIGHT)
add_slide_header(slide4, "02 AI 幫我們看的 4 個 ESG 檢查點")

# 2x2 grid of cards representing the 4 check points without classmate's or AI images
bullets4_1 = [
    ("• 任務目標：", "判定文本段落中是否包含企業具體的 ESG 承諾。"),
    ("• 任務設計：", "二分類 (是/否) | 任務權重 20%")
]
bullets4_2 = [
    ("• 任務目標：", "辨識該承諾規劃在何時兌現與實現。"),
    ("• 任務設計：", "五分類 (已完成/2年內/2-5年內/5年以上/NA) | 任務權重 15%")
]
bullets4_3 = [
    ("• 任務目標：", "分析文本是否包含支持該承諾的具體數據或實績。"),
    ("• 任務設計：", "三分類 (是/否/不適用) | 任務權重 30%")
]
bullets4_4 = [
    ("• 任務目標：", "評估所提佐證證據的可信度、清晰度與真實性。"),
    ("• 任務設計：", "四分類 (清晰/不清晰/具誤導性/NA) | 任務權重 35%")
]

add_card_text_content(slide4, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.2), "1. 承諾狀態 (promise_status)", bullets4_1, PRIMARY, font_size=13)
add_card_text_content(slide4, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.2), "2. 驗證時程 (verification_timeline)", bullets4_2, SECONDARY, font_size=13)
add_card_text_content(slide4, Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.3), "3. 證據狀態 (evidence_status)", bullets4_3, SECONDARY, font_size=13)
add_card_text_content(slide4, Inches(6.8), Inches(4.3), Inches(5.7), Inches(2.3), "4. 證據品質 (evidence_quality)", bullets4_4, PRIMARY, font_size=13)

# ------------------------------------------------------------
# Slide 5: 相關研究 - AI 是怎麼學習的？ (BERT Model Concept)
# ------------------------------------------------------------
slide5 = prs.slides.add_slide(slide_layout)
set_background(slide5, BG_LIGHT)
add_slide_header(slide5, "03 AI 是怎麼學習的？(BERT 學習原理)")

# Two side-by-side cards for Pre-training and Fine-tuning
bullets5_left = [
    ("圖書館廣泛閱讀：", "AI 大腦先在網路上閱讀海量無標籤的中文文章與書籍，學習中文的基本語意、文法、造句規律。這階段叫「預訓練」。"),
    ("全詞遮罩 (WWM) 技術：", "不只遮住單一字元，而是把完整的詞彙（如「碳中和」）遮起來讓 AI 猜，使其更懂中文實體與詞意。")
]
add_card_text_content(slide5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "第一步：預訓練 (學會讀書)", bullets5_left, PRIMARY)

bullets5_right = [
    ("針對考卷重點複習：", "一旦學會讀書後，我們在 AI 上方加裝簡單的預測層，使用特定任務的標記資料（如 ESG 考卷）進行訓練。這階段叫「微調」。"),
    ("一腦多用 (Multi-Task)：", "讓同一個 AI 大腦同時微調四個子任務。不僅能節省硬體計算資源，還能共享不同問題間的邏輯關聯，提升效果。")
]
add_card_text_content(slide5, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), "第二步：微調 (重點複習與答題)", bullets5_right, SECONDARY)

# ------------------------------------------------------------
# Slide 6: 研究方法 - 文本長度分布 (Text Length Distribution)
# ------------------------------------------------------------
slide6 = prs.slides.add_slide(slide_layout)
set_background(slide6, BG_LIGHT)
add_slide_header(slide6, "04 研究方法 - 文本長度分佈")

bullets6 = [
    ("長度集中區間：", "統計數據集中的 1,000 筆永續報告段落，發現大部份的段落長度集中在 300 到 400 字之間。"),
    ("長句的比例：", "字數超過 512 的長句在資料集中非常罕見，合計只有 17 筆，僅佔整體的 1.7%。"),
    ("最大長度設定：", "將 AI 的閱讀上限（MAX_LEN）設為 512，即可完整覆蓋 98.3% 的文本段落。這保證了長句後半段的細節和證據不被截斷，能有效拉升評分。")
]
add_card_text_content(slide6, Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.8), "長度統計與長句覆蓋", bullets6, SECONDARY)

# Right panel: Insert actual length distribution chart from PDF
add_slide_picture(slide6, "text_length_distribution.png", Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8), "文本長度分布圖")

# ------------------------------------------------------------
# Slide 7: 研究方法 - Tokenizer 左截斷的妙用 (Left Truncation)
# ------------------------------------------------------------
slide7 = prs.slides.add_slide(slide_layout)
set_background(slide7, BG_LIGHT)
add_slide_header(slide7, "05 研究方法 - 左截斷 (Left Truncation) 的妙用")

# Q&A format using two large side-by-side cards
bullets7_left = [
    ("預設是右截斷 (砍尾巴)：", "當段落字數超過 512 時，預設的 truncation 會直接把超過部分的「結尾」砍掉，保留開頭。"),
    ("為什麼砍尾巴很危險？", "在 ESG 的長文本中，開頭通常是大量的商業鋪陳、背景介紹和宏大的理想宣告，並不包含實質檢驗特徵。")
]
add_card_text_content(slide7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "【問】為什麼不砍後半部？", bullets7_left, PRIMARY)

bullets7_right = [
    ("改為左截斷 (砍頭留尾)：", "我們在設定 Tokenizer 時，將截斷邊調整為 left。這能讓模型在字數超限時捨棄開頭，保留末尾的 512 字。"),
    ("後半部通常寫些什麼？", "企業具體的「承諾兌現時程（時間線）」以及「佐證數據實績（證據）」，幾乎都出現在段落的最後面！"),
    ("結果：", "改為左截斷後，完美保留了關鍵特徵，使證據品質與時程判定分數顯著攀升。")
]
add_card_text_content(slide7, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), "【答】因為具體證據都在結尾！", bullets7_right, SECONDARY)

# ------------------------------------------------------------
# Slide 8: 研究方法 - 技術 1：特徵二合一融合 (Feature Fusion)
# ------------------------------------------------------------
slide8 = prs.slides.add_slide(slide_layout)
set_background(slide8, BG_LIGHT)
add_slide_header(slide8, "06 研究方法 - 特徵二合一融合 (Feature Fusion)")

bullets8_left = [
    ("• [CLS] 全局特徵：", "提取句子開頭的 [CLS] 向量（768維），代表文本段落的全局核心意圖。"),
    ("• Mean Pooling 特徵：", "對所有 token 進行平均池化（768維），代表整段文本的平均語意分佈。"),
    ("• 特徵拼接 (Concat)：", "拼接為 1536 維融合向量，同時捕捉全局核心與細節詞彙。")
]
bullets8_right = [
    ("• 共享 Backbone 網路：", "使用預訓練 RoBERTa-base 作為共享翻譯器，大幅節省運算資源。"),
    ("• 獨立雙層 MLP 頭：", "為 4 個任務設計獨立的 MLP 分類頭（加 Dropout），預防梯度干擾。"),
    ("• 協同多任務學習：", "多任務共同訓練、共享表徵，讓不同問題間的語意邏輯相互促進。")
]
add_card_text_content(slide8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "特徵二合一融合技術", bullets8_left, SECONDARY)
add_card_text_content(slide8, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), "多任務 MLP 分類架構", bullets8_right, PRIMARY)

# ------------------------------------------------------------
# Slide 9: 研究方法 - 技術 2：特例加權與決策門檻調整 (Loss & Threshold)
# ------------------------------------------------------------
slide9 = prs.slides.add_slide(slide_layout)
set_background(slide9, BG_LIGHT)
add_slide_header(slide9, "07 研究方法 - 特例加權與決策門檻調整")

bullets9_left = [
    ("資料不均衡問題：", "資料集裡「誤導性證據（Misleading）」或「兩年內兌現（within_2_years）」等特例樣本非常罕見（不到 5%）。"),
    ("逆頻率損失權重：", "如果 AI 猜錯了這些稀有類別，我們會對其施加數倍的懲罰，強制大腦花更多心思來強記和學習這些稀有樣本。")
]
add_card_text_content(slide9, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "優化 1：特例加權 (Inverse-Freq Loss)", bullets9_left, PRIMARY)

bullets9_right = [
    ("不要死猜概率最大的答案：", "常規的預測會直接取最大機率的值（argmax），但在類別極度不均時，AI 會本能地傾向猜常見標籤。"),
    ("決策信心門檻調整：", "我們透過網格搜索微調了少數類別的判定機率乘數，手動調低判定門檻。這極大地提高了少數類別的召回率，進一步拉升了 Macro F1 分數。")
]
add_card_text_content(slide9, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), "優化 2：決策信心門檻調整", bullets9_right, SECONDARY)

# ------------------------------------------------------------
# Slide 10: 研究方法 - 技術 3：三個臭皮匠與顯卡減重 (Ensemble & VRAM)
# ------------------------------------------------------------
slide10 = prs.slides.add_slide(slide_layout)
set_background(slide10, BG_LIGHT)
add_slide_header(slide10, "08 研究方法 - 三個臭皮匠與顯卡減重")

# Three vertical columns for clean, elegant, text-only card structure (no AI images)
bullets10_left = [
    ("• 消除偶發偏見：", "單一模型對特定資料易偏誤，集成預測消除偶然失誤。"),
    ("• 概率平均投票：", "將 3 個 Fold 獨立模型之預測概率平均，決策更穩健。")
]
bullets10_mid = [
    ("• 16位元半精度：", "使用 FP16 半精度進行混合精度計算，運算速度提升。"),
    ("• 記憶體減半：", "有效減少顯存消耗 50% 以上，保持模型效能與精準。")
]
bullets10_right = [
    ("• 本機顯存控制：", "批次大小設為 8，使本機顯存佔用控制在 3.0GB 以內。"),
    ("• 杜絕溢出當機：", "家用 4GB 顯卡 (RTX 3050) 也能順暢運行，防止 OOM。")
]

add_card_text_content(slide10, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "3折集成投票 (Ensemble)", bullets10_left, PRIMARY, font_size=13)
add_card_text_content(slide10, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "自動混合精度 (AMP)", bullets10_mid, SECONDARY, font_size=13)
add_card_text_content(slide10, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8), "批次與顯存優化", bullets10_right, PRIMARY, font_size=13)

# ------------------------------------------------------------
# Slide 11: 實驗設定 - 資料集與 3-Fold 切分 (Datasets & CV)
# ------------------------------------------------------------
slide11 = prs.slides.add_slide(slide_layout)
set_background(slide11, BG_LIGHT)
add_slide_header(slide11, "09 實驗設定 - 資料集與 3-Fold 切分")

bullets11_left = [
    ("資料集來源：", "採用 AICUP 競賽提供的 `vpesg4k_train_1000.json`，共 1,000 筆有標籤的永續報告中文文本段落。"),
    ("本機驗證方式：", "採用 3 折交叉驗證 (3-Fold Cross-Validation) 作為我們本地開發與模型評估的唯一準則。"),
    ("重複性設定：", "固定亂數種子 `random_state=42`，確保每次資料切割和模型訓練都可以百分之百還原，保證實驗可重複。")
]
add_card_text_content(slide11, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "評估策略與資料集", bullets11_left, SECONDARY)

bullets11_right = [
    ("為什麼要用 3 折交叉驗證？", "如果只用普通的 2/8 分法，結果很容易受到特定劃分影響。"),
    ("3-Fold 的運作方式：", "將 1000 筆隨機分為 A、B、C 三份。\n- 輪流用 (A+B) 讀書，C 考試；\n- (B+C) 讀書，A 考試；\n- (A+C) 讀書，B 考試。\n- 最終綜合 3 次考卷的平均成績。這能反映 AI 真正的實力，不會偏袒任何特定資料。")
]
add_card_text_content(slide11, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), "什麼是 3 折交叉驗證？", bullets11_right, PRIMARY)

# ------------------------------------------------------------
# Slide 12: 實驗設定 - 四大子任務分布不均 (Task Imbalance)
# ------------------------------------------------------------
slide12 = prs.slides.add_slide(slide_layout)
set_background(slide12, BG_LIGHT)
add_slide_header(slide12, "10 實驗設定 - 四大子任務分布不均")

# Two sections: Left card, Right detailed distribution table
bullets12 = [
    ("極端不平衡的分布：", "在 ESG 報告中，多數標籤都是偏向正面、清晰的（如承諾狀態為 Yes 佔 80%，證據品質 Clear 佔 85%）。"),
    ("少數類別的生存危機：", "少數類別的樣本非常稀缺。例如，在 1000 筆中，「具誤導性的證據（Misleading）」僅有幾十個樣本，如果使用常規訓練，AI 會直接放棄學習牠們。"),
    ("優化對策：", "這也是我們必須導入「特例加權」與「決策信心門檻調整」的核心原因。")
]
add_card_text_content(slide12, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "資料失衡挑戰", bullets12, SECONDARY)

# Insert esg_label_distribution.png on the right as visual evidence of label imbalance
add_slide_picture(slide12, "esg_label_distribution.png", Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), "四大任務標籤分布不均分析")

# ------------------------------------------------------------
# Slide 13: 實驗結果 - 模型分數大車拼 (Results Comparison Table)
# ------------------------------------------------------------
slide13 = prs.slides.add_slide(slide_layout)
set_background(slide13, BG_LIGHT)
add_slide_header(slide13, "11 實驗結果 - 模型分數大車拼")

rows = 5
cols = 6
table_left = Inches(0.8)
table_top = Inches(1.8)
table_width = Inches(6.5)
table_height = Inches(4.8)

table_shape = slide13.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
table = table_shape.table

table.columns[0].width = Inches(1.8)  # Model name
table.columns[1].width = Inches(1.0)  # Weighted F1
table.columns[2].width = Inches(0.9)  # promise
table.columns[3].width = Inches(0.9)  # timeline
table.columns[4].width = Inches(0.9)  # evidence
table.columns[5].width = Inches(1.0)  # quality

headers = ["模型對比", "綜合得分", "承諾 (20%)", "時程 (15%)", "證據 (30%)", "品質 (35%)"]
data = [
    ["1. Baseline (一般 BERT)", "0.56067", "0.75685", "0.43928", "0.60017", "0.46674"],
    ["2. 本機優化單模型", "0.59016", "0.75758", "0.48172", "0.64089", "0.49747"],
    ["3. 本機最優集成 (Tuned)", "0.59794", "0.79266", "0.50429", "0.67344", "0.46211"],
    ["4. 雲端大模型集成版", "0.76801", "0.86542", "0.69874", "0.79512", "0.73248"]
]

for col_idx, text in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.name = FONT_NAME
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

for row_idx, row_data in enumerate(data):
    for col_idx, text in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.fill.solid()
        if row_idx == 2:
            cell.fill.fore_color.rgb = RGBColor(255, 243, 215) # Soft gold
        elif row_idx == 3:
            cell.fill.fore_color.rgb = RGBColor(220, 240, 230) # Soft green
        else:
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 0 else RGBColor(240, 242, 245)
            
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.name = FONT_NAME
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK
        p.alignment = PP_ALIGN.CENTER
        if col_idx == 1 or row_idx >= 2:
            p.font.bold = True

# Insert model_comparison.png on the right as visual comparison
add_slide_picture(slide13, "model_comparison.png", Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.8), "模型加權 F1 得分對比")

# ------------------------------------------------------------
# Slide 14: 實驗結果 - 四大任務預測表現 (Task Performance)
# ------------------------------------------------------------
slide14 = prs.slides.add_slide(slide_layout)
set_background(slide14, BG_LIGHT)
add_slide_header(slide14, "12 實驗結果 - 四大子任務預測表現")

bullets14 = [
    ("承諾狀態 (F1: 0.8654)：", "表現最突出。AI 能夠非常精準地看穿這段文字究竟是個空泛的理想宣示，還是真正提出了具體的承諾目標。"),
    ("證據狀態 (F1: 0.7951)：", "表現優異。AI 能夠正確判斷句子中是否提到了與承諾對應的實質數字或佐證事實。"),
    ("證據品質 (F1: 0.7324)：", "藉由逆頻率損失和門檻調優，AI 能抓出那些「模糊、誤導性」的文字，評分顯著拉升。"),
    ("驗證時程 (F1: 0.6987)：", "因為時程有 5 個類別（多分類任務），難度最高，但我們的方法依然比 Baseline (0.439) 提升了 59%。")
]
add_card_text_content(slide14, Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.8), "四大檢查任務的精準度分析 (雲端集成版)", bullets14, SECONDARY, font_size=14)

# Right panel: Insert user's original subtask performance chart
add_slide_picture(slide14, "subtask_performance.png", Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8), "子任務表現分析圖")

# ------------------------------------------------------------
# Slide 15: 實驗結果 - 訓練曲線與特徵學習 (Visual Evidence)
# ------------------------------------------------------------
slide15 = prs.slides.add_slide(slide_layout)
set_background(slide15, BG_LIGHT)
add_slide_header(slide15, "13 實驗結果 - 訓練曲線與特徵學習 (圖表佐證)")

bullets15 = [
    ("Loss 穩定下降：", "從隨附的訓練曲線圖可以看出，在訓練過程中，Loss (損失值) 隨著學習進程呈現平滑且穩定的下滑，代表模型正在有效地提取特徵，沒有出現明顯的梯度震盪。"),
    ("F1 分數穩定收斂：", "加權 F1 分數（紅色虛線）在第 8 個 Epoch 達到最高點後即趨於穩定，說明模型收斂平穩，沒有過擬合 (Overfitting) 的現象。"),
    ("客觀圖表實證：", "這張圖表為我們的特徵融合與多任務學習模式提供了堅實的客觀數據支持，證實該優化策略極具實用性。")
]
add_card_text_content(slide15, Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.8), "模型學習曲線佐證", bullets15, PRIMARY)
add_slide_picture(slide15, "optimized_curve.png", Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8), "模型學習曲線")

# ------------------------------------------------------------
# Slide 16: 結論與未來展望 (Conclusion & Future Work)
# ------------------------------------------------------------
slide16 = prs.slides.add_slide(slide_layout)
set_background(slide16, BG_LIGHT)
add_slide_header(slide16, "14 結論與未來展望")

bullets16_left = [
    ("• 自動化查核系統：", "成功開發基於 RoBERTa 的多任務 ESG 驗證系統，能客觀、高效辨識企業「漂綠」風險。"),
    ("• 優化策略得到印證：", "特徵融合、多模型集成與決策門檻調優，將加權 F1 顯著提升 36.98%，證實本研究框架之有效性。")
]
bullets16_right = [
    ("• 本土化語意增強：", "規劃擴大蒐集台灣上市櫃企業永續報告書，進行二次預訓練，深化本土永續術語之語境理解。"),
    ("• 推理確信與LLM整合：", "計畫結合生成式大語言模型（LLM）之上下文推理能力，以應對更複雜的 ESG 指標確信查核任務。")
]
add_card_text_content(slide16, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "主要結論 (Conclusion)", bullets16_left, PRIMARY)
add_card_text_content(slide16, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), "未來展望 (Future Outlook)", bullets16_right, SECONDARY)

# ------------------------------------------------------------
# Slide 17: Ending Slide (Dark Theme)
# ------------------------------------------------------------
slide17 = prs.slides.add_slide(slide_layout)
set_background(slide17, BG_DARK)

# Background Accent Shapes
acc1_end = slide17.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(3.833), Inches(7.5))
acc1_end.fill.solid()
acc1_end.fill.fore_color.rgb = RGBColor(24, 38, 60)
remove_border(acc1_end)

acc2_end = slide17.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(3.833), Inches(0), Inches(2.0), Inches(7.5))
acc2_end.fill.solid()
acc2_end.fill.fore_color.rgb = RGBColor(24, 38, 60)
remove_border(acc2_end)

acc3_end = slide17.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(7.5))
acc3_end.fill.solid()
acc3_end.fill.fore_color.rgb = SECONDARY
remove_border(acc3_end)

# Thank you text box
thank_box = slide17.shapes.add_textbox(Inches(4.5), Inches(2.3), Inches(7.5), Inches(2.5))
tf_thank = thank_box.text_frame
tf_thank.word_wrap = True
tf_thank.margin_left = tf_thank.margin_right = tf_thank.margin_top = tf_thank.margin_bottom = 0

pt1 = tf_thank.paragraphs[0]
pt1.text = "簡報結束，感謝聆聽"
pt1.font.name = FONT_NAME
pt1.font.size = Pt(40)
pt1.font.bold = True
pt1.font.color.rgb = WHITE

pt2 = tf_thank.add_paragraph()
pt2.text = "讓 ESG 承諾與佐證證據查核更簡單、更直覺、更可信"
pt2.font.name = FONT_NAME
pt2.font.size = Pt(16)
pt2.font.color.rgb = SECONDARY
pt2.space_before = Pt(14)

pt3 = tf_thank.add_paragraph()
pt3.text = "歡迎指教與討論"
pt3.font.name = FONT_NAME
pt3.font.size = Pt(18)
pt3.font.color.rgb = ACCENT
pt3.space_before = Pt(8)

# ------------------------------------------------------------
# Save Presentation
# ------------------------------------------------------------
output_filename = "ESG_Promise_Verification_Presentation.pptx"
prs.save(output_filename)
print(f"Presentation successfully saved to: {os.path.abspath(output_filename)}")
